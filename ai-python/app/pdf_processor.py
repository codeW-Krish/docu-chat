import pdfplumber
import pytesseract
import os
import docx2txt
import pandas as pd
import csv
from pptx import Presentation

# Configure Tesseract path from environment
pytesseract.pytesseract.tesseract_cmd = os.getenv('TESSERACT_PATH', 'tesseract')
from PIL import Image
import io
import os
import uuid
import logging
from langchain.text_splitter import RecursiveCharacterTextSplitter

from .embedding_service import EmbeddingService
from .database import get_db_connection

logger = logging.getLogger(__name__)

class DocumentProcessor:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,      # Increased chunk size for better context
            chunk_overlap=100,   # Increased overlap for better context
            length_function=len,
            separators=["\n\n", "\n", ". ", "! ", "? ", " ", ""]
        )
        self.embedding_service = EmbeddingService()
    
    def extract_text_from_docx(self, file_path):
        """Extract text from DOCX file"""
        try:
            text = docx2txt.process(file_path)
            return [{'page_number': 1, 'text': text, 'char_length': len(text), 'used_ocr': False}]
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}")
            raise

    def extract_text_from_txt(self, file_path):
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            return [{'page_number': 1, 'text': text, 'char_length': len(text), 'used_ocr': False}]
        except Exception as e:
            logger.error(f"TXT extraction error: {str(e)}")
            raise

    def extract_text_from_csv(self, file_path):
        """Extract text from CSV file"""
        try:
            df = pd.read_csv(file_path)
            text = df.to_string(index=False)
            return [{'page_number': 1, 'text': text, 'char_length': len(text), 'used_ocr': False}]
        except Exception as e:
            logger.error(f"CSV extraction error: {str(e)}")
            raise

    def extract_text_from_pptx(self, file_path):
        """Extract text from PPTX file (slides and notes)"""
        try:
            prs = Presentation(file_path)
            pages_data = []
            
            for i, slide in enumerate(prs.slides):
                text_content = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                
                # Extract text from notes if available
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes:
                        text_content.append(f"\n[Notes]: {notes}")
                
                full_text = "\n".join(text_content)
                
                pages_data.append({
                    'page_number': i + 1,
                    'text': full_text,
                    'char_length': len(full_text),
                    'used_ocr': False
                })
                
            return pages_data
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}")
            raise

    def extract_text_with_ocr(self, pdf_path):
        """Extract text from PDF with intelligent OCR fallback"""
        pages_data = []
        
        try:
            with pdfplumber.open(pdf_path) as pdf:
                total_pages = len(pdf.pages)
                logger.info(f"Processing PDF with {total_pages} pages")
                
                for page_num, page in enumerate(pdf.pages, 1):
                    logger.info(f"Processing page {page_num}/{total_pages}")
                    
                    # Primary text extraction using pdfplumber
                    page_text = page.extract_text() or ""
                    has_sufficient_text = len(page_text.strip()) > 100
                    
                    ocr_text = ""
                    used_ocr = False
                    
                    # OCR fallback for image-heavy pages
                    if not has_sufficient_text:
                        logger.info(f"Low text on page {page_num}, attempting OCR...")
                        try:
                            # Convert page to high-resolution image
                            page_image = page.to_image(resolution=300)
                            
                            # Convert to PIL Image for OCR
                            img_bytes = io.BytesIO()
                            page_image.save(img_bytes, format='PNG', quality=100)
                            img_bytes.seek(0)
                            
                            # Perform OCR with optimized settings
                            ocr_text = pytesseract.image_to_string(
                                Image.open(img_bytes), 
                                config='--psm 6 -c preserve_interword_spaces=1'
                            )
                            used_ocr = bool(ocr_text.strip())
                            
                            if used_ocr:
                                logger.info(f"OCR extracted {len(ocr_text)} characters from page {page_num}")
                            else:
                                logger.warning(f"OCR failed to extract text from page {page_num}")
                                
                        except Exception as ocr_error:
                            logger.warning(f"OCR failed for page {page_num}: {str(ocr_error)}")
                    
                    # Combine text sources
                    final_text = page_text.strip()
                    if used_ocr and ocr_text.strip():
                        if final_text:
                            final_text += "\n\n" + ocr_text.strip()
                        else:
                            final_text = ocr_text.strip()
                    
                    pages_data.append({
                        'page_number': page_num,
                        'text': final_text,
                        'char_length': len(final_text),
                        'used_ocr': used_ocr,
                        'has_native_text': bool(page_text.strip())
                    })
                    
                    logger.info(f"Page {page_num}: {len(final_text)} chars (OCR: {used_ocr})")
        
        except Exception as e:
            logger.error(f"PDF processing error: {str(e)}")
            raise
        
        return pages_data
    
    def chunk_text_with_metadata(self, text, page_number):
        """Split text into chunks with proper character offsets and metadata"""
        if not text.strip():
            return []
            
        chunks = self.text_splitter.split_text(text)
        chunk_objects = []
        
        current_char = 0
        for chunk_index, chunk in enumerate(chunks):
            if not chunk.strip():
                continue
                
            start_char = current_char
            end_char = current_char + len(chunk)
            
            chunk_objects.append({
                'chunk_index': chunk_index,
                'page_number': page_number,
                'start_char': start_char,
                'end_char': end_char,
                'chunk_text': chunk.strip(),
                'word_count': len(chunk.split())
            })
            
            # Move position, accounting for overlap
            current_char = end_char - 50  # Approximate overlap adjustment
        
        logger.info(f"Page {page_number} split into {len(chunk_objects)} chunks")
        return chunk_objects
    
    def process_pdf(self, pdf_id, file_path, user_id):
        import tempfile
        from appwrite.client import Client
        from appwrite.services.storage import Storage
        
        tmp_file_path = None
        conn = None
        try:
            logger.info(f"Starting document processing: {pdf_id} for user {user_id}")
            
            # Initialize Appwrite client
            client = Client()
            client.set_endpoint(os.getenv('APPWRITE_ENDPOINT'))
            client.set_project(os.getenv('APPWRITE_PROJECT_ID'))
            client.set_key(os.getenv('APPWRITE_API_KEY'))
            
            storage = Storage(client)
            bucket_id = os.getenv('APPWRITE_STORAGE_BUCKET_ID')
            
            # Here file_path actually contains the Appwrite File ID from PHP backend
            file_id = file_path 
            logger.info(f"Downloading file ID {file_id} from Appwrite bucket {bucket_id}")
            
            try:
                # Bypass the buggy SDK get_file_download by using the REST API directly
                import requests
                
                endpoint = os.getenv('APPWRITE_ENDPOINT')
                project = os.getenv('APPWRITE_PROJECT_ID')
                api_key = os.getenv('APPWRITE_API_KEY')
                
                url = f"{endpoint}/storage/buckets/{bucket_id}/files/{file_id}/download"
                headers = {
                    "X-Appwrite-Project": project,
                    "X-Appwrite-Key": api_key
                }
                
                logger.info(f"Downloading directly from Appwrite API: {url}")
                response = requests.get(url, headers=headers)
                
                if response.status_code != 200:
                    raise Exception(f"Appwrite API returned {response.status_code}: {response.text}")
                    
                result = response.content
                
                # Try to get the file name/extension from Appwrite metadata, fallback to .pdf
                try:
                    meta = storage.get_file(bucket_id, file_id)
                    original_name = meta['name']
                    ext = os.path.splitext(original_name)[1].lower()
                    if not ext:
                        ext = '.pdf'
                except Exception as meta_ex:
                    logger.warning(f"Failed to get file metadata, defaulting to .pdf: {meta_ex}")
                    ext = '.pdf'
                
                # Save to a temporary file locally
                fd, tmp_file_path = tempfile.mkstemp(suffix=ext)
                with os.fdopen(fd, 'wb') as f:
                    f.write(result)
                    
            except Exception as e:
                logger.error(f"Failed to download file from Appwrite: {e}")
                raise FileNotFoundError(f"File ID {file_id} could not be downloaded from Appwrite")
            
            # Determine file type and extract text using the temp file
            logger.info(f"Processing temporary file: {tmp_file_path}, detected extension: {ext}")
            if ext == '.pdf':
                pages_data = self.extract_text_with_ocr(tmp_file_path)
            elif ext == '.docx':
                pages_data = self.extract_text_from_docx(tmp_file_path)
            elif ext == '.txt':
                pages_data = self.extract_text_from_txt(tmp_file_path)
            elif ext == '.csv':
                pages_data = self.extract_text_from_csv(tmp_file_path)
            elif ext == '.pptx':
                pages_data = self.extract_text_from_pptx(tmp_file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
            
            # Check if we extracted any content
            total_text_length = sum(page['char_length'] for page in pages_data)
            if total_text_length == 0:
                raise ValueError("No text content could be extracted from document")
                
            total_pages = len(pages_data)
            
            logger.info(f"Extracted {total_text_length} characters from {total_pages} pages/sections")
            
            # Open a secondary connection strictly for live UI telemetry
            telemetry_conn = get_db_connection()
            telemetry_conn.autocommit = True
            try:
                with telemetry_conn.cursor() as t_cursor:
                    t_cursor.execute(
                        "UPDATE pdfs SET processing_progress = 30, page_count = %s WHERE pdf_id = %s", 
                        (total_pages, pdf_id)
                    )
            except Exception as e:
                logger.warning(f"Telemetry update failed: {e}")
                
            # Get database connection for main transaction
            conn = get_db_connection()
            cursor = conn.cursor()
            
            # Start transaction
            conn.autocommit = False
            
            all_chunks = []
            
            # Process each page and chunk text
            for page_data in pages_data:
                if page_data['text']:
                    page_chunks = self.chunk_text_with_metadata(
                        page_data['text'], 
                        page_data['page_number']
                    )
                    all_chunks.extend(page_chunks)
            
            logger.info(f"Total chunks created: {len(all_chunks)}")
            
            # Insert chunks and generate embeddings in transaction
            successful_chunks = 0
            for chunk_data in all_chunks:
                try:
                    chunk_id = str(uuid.uuid4())
                    
                    # Insert into pdf_chunks (metadata)
                    cursor.execute("""
                        INSERT INTO pdf_chunks 
                        (chunk_id, pdf_id, user_id, chunk_index, page_number, start_char, end_char, chunk_text)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                    """, (
                        chunk_id, pdf_id, user_id, chunk_data['chunk_index'],
                        chunk_data['page_number'], chunk_data['start_char'],
                        chunk_data['end_char'], chunk_data['chunk_text']
                    ))
                    
                    # Generate embedding
                    embedding = self.embedding_service.generate_embedding(chunk_data['chunk_text'])
                    
                    # Insert into pdf_chunks_embeddings
                    cursor.execute("""
                        INSERT INTO pdf_chunks_embeddings 
                        (chunk_id, pdf_id, user_id, chunk_index, chunk_text, start_char, end_char, embedding)
                        VALUES (%s::uuid, %s::uuid, %s::uuid, %s, %s, %s, %s, %s::vector)
                    """, (
                        chunk_id, pdf_id, user_id, chunk_data['chunk_index'],
                        chunk_data['chunk_text'], chunk_data['start_char'],
                        chunk_data['end_char'], embedding
                    ))
                    
                    successful_chunks += 1
                    
                except Exception as chunk_error:
                    logger.error(f"Failed to process chunk {chunk_data['chunk_index']}: {str(chunk_error)}")
                    # Continue with other chunks
                    continue
                    
                # Calculate real-time percentage (30% to 95%) and stream it via telemetry
                if successful_chunks % 5 == 0 or successful_chunks == len(all_chunks):
                    # Progress bounds: 30% base + up to 65% for vectorization
                    current_progress = 30 + int((successful_chunks / max(1, len(all_chunks))) * 65)
                    try:
                        with telemetry_conn.cursor() as t_cursor:
                            t_cursor.execute("UPDATE pdfs SET processing_progress = %s WHERE pdf_id = %s", (current_progress, pdf_id))
                    except Exception as e:
                        pass
            
            # Update PDF processing status to 100% completed
            cursor.execute("""
                UPDATE pdfs 
                SET processing_status = 'completed', processing_progress = 100, page_count = %s
                WHERE pdf_id = %s
            """, (len(pages_data), pdf_id))
            
            # Commit transaction
            conn.commit()
            
            logger.info(f"Document processing completed: {pdf_id}")
            logger.info(f"Statistics: {successful_chunks}/{len(all_chunks)} chunks processed")
            
            return {
                'status': 'success',
                'message': 'Document processed successfully',
                'data': {
                    'chunk_count': successful_chunks,
                    'page_count': len(pages_data),
                    'total_chunks': len(all_chunks)
                }
            }
            
        except Exception as e:
            # Rollback transaction on error
            if conn:
                conn.rollback()
                logger.error("Transaction rolled back due to error")
            
            logger.error(f"Document processing failed: {str(e)}")
            
            # Update PDF status to error
            if conn:
                try:
                    cursor.execute("UPDATE pdfs SET processing_status = 'failed' WHERE pdf_id = %s", (pdf_id,))
                    conn.commit()
                except Exception as update_error:
                    logger.error(f"Failed to update document status: {str(update_error)}")
            
            return {
                'status': 'error',
                'message': f'Document processing failed: {str(e)}'
            }
            
        finally:
            # Cleanup temporary file
            if tmp_file_path and os.path.exists(tmp_file_path):
                try:
                    os.remove(tmp_file_path)
                    logger.info(f"Cleaned up temporary file: {tmp_file_path}")
                except Exception as e:
                    logger.error(f"Failed to cleanup temp file {tmp_file_path}: {e}")
                    
            # Close connection
            if conn:
                conn.close()
                
            if 'telemetry_conn' in locals() and telemetry_conn:
                telemetry_conn.close()
